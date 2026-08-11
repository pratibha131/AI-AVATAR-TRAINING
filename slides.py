"""
Slide engine: pptx/pdf upload -> exact-fidelity slide images + progressive
"build state" images (bullet-by-bullet reveals) rendered through LibreOffice
so the original styling is preserved pixel-for-pixel.

Build states are produced by duplicating each slide inside the pptx zip and
making not-yet-revealed paragraphs fully transparent (alpha 0) — layout never
reflows, so every state is the exact original slide with later bullets hidden.
"""
import os, re, math, shutil, subprocess, zipfile, copy, tempfile, uuid, json
from lxml import etree
from pptx import Presentation
from pptx.util import Emu

NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'ct': 'http://schemas.openxmlformats.org/package/2006/content-types',
    'rel': 'http://schemas.openxmlformats.org/package/2006/relationships',
}
A = NS['a']; P = NS['p']; R = NS['r']

MAX_STEPS_PER_SLIDE = 25
BUILD_VERSION = 3   # bump when the ordering/segmentation logic changes so
                    # stale pre-upgrade build states get rebuilt automatically


def _q(tag):  # a:foo -> qualified
    pfx, local = tag.split(':')
    return '{%s}%s' % (NS[pfx], local)


def _run_soffice(args, timeout=300):
    env = dict(os.environ)
    for attempt in range(3):
        try:
            subprocess.run(['soffice', '--headless', *args], check=True,
                           timeout=timeout, capture_output=True, env=env)
            return
        except subprocess.CalledProcessError:
            # kill a wedged instance between retries (platform-appropriate,
            # best-effort — a failed kill must not abort the retry loop)
            try:
                if os.name == 'nt':
                    for image in ('soffice.bin', 'soffice.exe'):
                        subprocess.run(['taskkill', '/F', '/IM', image],
                                       capture_output=True)
                else:
                    subprocess.run(['pkill', '-f', 'soffice'],
                                   capture_output=True)
            except Exception:
                pass
            import time; time.sleep(2)
    raise RuntimeError('LibreOffice conversion failed')


def _num_sort_key(filename):
    m = re.search(r'\d+', filename)
    return int(m.group()) if m else 0


def _pptx_slide_count(pptx_path):
    """Number of slides straight from the package (cheap, no Office needed)."""
    try:
        with zipfile.ZipFile(pptx_path) as z:
            root = etree.fromstring(z.read('ppt/presentation.xml'))
        lst = root.find(_q('p:sldIdLst'))
        return len(lst) if lst is not None else 0
    except Exception:
        return None


def _export_via_powerpoint(pptx_path, outdir, prefix='slide'):
    """Native PowerPoint COM export (Windows). Returns PNG paths or None.

    Must work from worker threads: COM requires per-thread initialization,
    and DispatchEx keeps us in our own PowerPoint instance so we never
    close a presentation the user has open."""
    try:
        import pythoncom, win32com.client
    except ImportError:
        return None
    pythoncom.CoInitialize()
    ppt = pres = None
    tmp_exp = None
    try:
        abs_pptx = os.path.abspath(pptx_path)
        abs_out = os.path.abspath(outdir)
        tmp_exp = os.path.join(abs_out, '_tmp_exp_' + uuid.uuid4().hex[:6])
        os.makedirs(tmp_exp, exist_ok=True)
        ppt = win32com.client.DispatchEx('PowerPoint.Application')
        pres = ppt.Presentations.Open(abs_pptx, WithWindow=False)
        n_slides = int(pres.Slides.Count)
        pres.SaveAs(tmp_exp, 18)  # 18 = ppSaveAsPNG
        exp_files = sorted((f for f in os.listdir(tmp_exp)
                            if f.lower().endswith(('.png', '.jpg'))),
                           key=_num_sort_key)
        if len(exp_files) != n_slides:
            raise RuntimeError(f'PowerPoint exported {len(exp_files)} images '
                               f'for {n_slides} slides')
        out_files = []
        for idx, fname in enumerate(exp_files):
            dst = os.path.join(abs_out, f"{prefix}-{idx+1:03d}.png")
            shutil.copy(os.path.join(tmp_exp, fname), dst)
            out_files.append(dst)
        return out_files or None
    except Exception as e:
        print(f'[slides] PowerPoint COM export failed: {e}')
        return None
    finally:
        try:
            if pres is not None:
                pres.Close()
        except Exception:
            pass
        try:
            if ppt is not None:
                ppt.Quit()
        except Exception:
            pass
        if tmp_exp:
            shutil.rmtree(tmp_exp, ignore_errors=True)
        pythoncom.CoUninitialize()


def pptx_to_pngs(pptx_path, outdir, dpi=110, prefix='slide', expected=None):
    """Export every slide to PNG. Guarantees the returned list matches the
    slide count (raises instead of silently dropping slides)."""
    os.makedirs(outdir, exist_ok=True)
    if expected is None:
        expected = _pptx_slide_count(pptx_path)

    out_files = _export_via_powerpoint(pptx_path, outdir, prefix)
    if out_files is None:
        # transient COM hiccups (busy instance) deserve one retry
        import time
        time.sleep(1.5)
        out_files = _export_via_powerpoint(pptx_path, outdir, prefix)

    if out_files is None:
        if shutil.which('soffice') is None or shutil.which('pdftoppm') is None:
            raise RuntimeError(
                'Slide export failed: PowerPoint COM export did not succeed '
                'and the LibreOffice/poppler fallback is incomplete (need '
                'both soffice and pdftoppm on PATH).')
        with tempfile.TemporaryDirectory() as td:
            _run_soffice(['--convert-to', 'pdf', pptx_path, '--outdir', td])
            pdfs = [f for f in os.listdir(td) if f.endswith('.pdf')]
            pdf = os.path.join(td, pdfs[0])
            subprocess.run(['pdftoppm', '-png', '-r', str(dpi), pdf,
                            os.path.join(outdir, prefix)], check=True, timeout=600)
        files = sorted((f for f in os.listdir(outdir)
                        if f.startswith(prefix + '-') and f.endswith('.png')),
                       key=_num_sort_key)
        out_files = [os.path.join(outdir, f) for f in files]

    if expected is not None and expected > 0 and len(out_files) != expected:
        raise RuntimeError(f'Slide export produced {len(out_files)} images '
                           f'but the deck has {expected} slides')
    return out_files


def pdf_to_pngs(pdf_path, outdir, dpi=110, prefix='slide'):
    os.makedirs(outdir, exist_ok=True)
    subprocess.run(['pdftoppm', '-png', '-r', str(dpi), pdf_path,
                    os.path.join(outdir, prefix)], check=True, timeout=600)
    files = sorted((f for f in os.listdir(outdir)
                    if f.startswith(prefix + '-') and f.endswith('.png')), key=_num_sort_key)
    return [os.path.join(outdir, f) for f in files]



# ---------------------------------------------------------------- extraction

def extract_slide_info(pptx_path):
    """Per slide: title, content paragraphs (in reading order), speaker notes."""
    prs = Presentation(pptx_path)
    slides = []
    for idx, slide in enumerate(prs.slides):
        title = ''
        paras = []            # list of {text, shape_id, para_idx}
        shapes_sorted = sorted(
            [sh for sh in slide.shapes if sh.has_text_frame],
            key=lambda sh: ((sh.top or 0), (sh.left or 0)))
        for sh in shapes_sorted:
            is_title = False
            try:
                if sh == slide.shapes.title:
                    is_title = True
            except Exception:
                pass
            for pi, para in enumerate(sh.text_frame.paragraphs):
                txt = ''.join(r.text for r in para.runs).strip()
                if not txt:
                    continue
                if is_title or (not title and (sh.top or 0) < prs.slide_height * 0.18):
                    if not title:
                        title = txt
                        continue
                paras.append({'text': txt, 'shape_id': sh.shape_id, 'para_idx': pi})
        notes = ''
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append({'index': idx, 'title': title, 'paragraphs': paras,
                       'notes': notes})
    return slides


def draft_script(slide_info):
    """Heuristic narration draft the user can edit (notes win if present)."""
    if slide_info['notes']:
        return slide_info['notes']
    parts = []
    if slide_info['title']:
        parts.append(f"Let's talk about {slide_info['title'].rstrip('?.! ')}.")
    for p in slide_info['paragraphs'][:8]:
        t = p['text'].rstrip()
        if len(t) < 3:
            continue
        if not re.search(r'[.!?]$', t):
            t += '.'
        parts.append(t)
    return ' '.join(parts) if parts else 'This slide speaks for itself.'


def draft_script_ordered(manifest_entry, slide_info=None):
    """Narration draft that walks the slide in its REVEAL (reading) order,
    so the spoken sequence and the element animations line up one-to-one.
    Speaker notes still win when present."""
    if slide_info and slide_info.get('notes'):
        return slide_info['notes']
    parts = []
    if slide_info and slide_info.get('title'):
        parts.append(f"Let's talk about {slide_info['title'].rstrip('?.! ')}.")
    seen = set()
    for t in (manifest_entry.get('texts') or []):
        t = re.sub(r'\s+', ' ', (t or '')).strip()
        if len(t) < 3 or t.lower() in seen:
            continue
        seen.add(t.lower())
        if not re.search(r'[.!?]$', t):
            t += '.'
        parts.append(t)
    return ' '.join(parts) if parts else (
        draft_script(slide_info) if slide_info else 'This slide speaks for itself.')


# ------------------------------------------------------- build-state engine

def _iter_body_paragraphs(sp_tree):
    """Yield (shape_el, txBody, [paragraph elements]) for non-title text shapes."""
    for sp in sp_tree.iter(_q('p:sp')):
        ph = sp.findall('.//' + _q('p:nvSpPr') + '/' + _q('p:nvPr') + '/' + _q('p:ph'))
        is_title = any(p.get('type') in ('title', 'ctrTitle') for p in ph)
        tx = sp.find(_q('p:txBody'))
        if tx is None or is_title:
            continue
        paras = tx.findall(_q('a:pPr') + '/..')  # not right; do directly
        yield sp, tx, tx.findall(_q('a:p'))


def _para_text(p_el):
    return ''.join(t.text or '' for t in p_el.iter(_q('a:t'))).strip()


def _make_transparent(p_el):
    """Make a paragraph invisible without changing layout."""
    # every run: force 0% alpha fill
    for r_el in list(p_el.iter(_q('a:r'))) + list(p_el.iter(_q('a:fld'))):
        rPr = r_el.find(_q('a:rPr'))
        if rPr is None:
            rPr = etree.SubElement(r_el, _q('a:rPr'))
            r_el.remove(rPr); r_el.insert(0, rPr)
        for fill in rPr.findall(_q('a:solidFill')):
            rPr.remove(fill)
        fill = etree.Element(_q('a:solidFill'))
        clr = etree.SubElement(fill, _q('a:srgbClr')); clr.set('val', 'FFFFFF')
        alpha = etree.SubElement(clr, _q('a:alpha')); alpha.set('val', '0')
        # attribute order: insert solidFill first child of rPr (after ln if any)
        rPr.insert(0, fill)
    # hide bullet glyph
    pPr = p_el.find(_q('a:pPr'))
    if pPr is None:
        pPr = etree.Element(_q('a:pPr'))
        p_el.insert(0, pPr)
    for tag in ('a:buChar', 'a:buAutoNum', 'a:buNone'):
        for e in pPr.findall(_q(tag)):
            pPr.remove(e)
    etree.SubElement(pPr, _q('a:buNone'))


def _hide_sp_visual(sp):
    """Make an entire shape invisible (fill, outline, effects, text) without
    touching its geometry — layout stays pixel-identical."""
    spPr = sp.find(_q('p:spPr'))
    if spPr is not None:
        for tag in ('a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill',
                    'a:grpFill', 'a:noFill', 'a:effectLst'):
            for e in spPr.findall(_q(tag)):
                spPr.remove(e)
        # explicit noFill overrides any style-inherited fill
        xfrm = spPr.find(_q('a:xfrm'))
        geom_idx = list(spPr).index(xfrm) + 1 if xfrm is not None else 0
        nf = etree.Element(_q('a:noFill'))
        # place after geometry elements (prstGeom/custGeom) if present
        insert_at = len(spPr)
        for i, ch in enumerate(spPr):
            if ch.tag in (_q('a:prstGeom'), _q('a:custGeom')):
                insert_at = i + 1
        spPr.insert(insert_at, nf)
        ln = spPr.find(_q('a:ln'))
        if ln is None:
            ln = etree.SubElement(spPr, _q('a:ln'))
        for tag in ('a:solidFill', 'a:gradFill', 'a:pattFill', 'a:noFill'):
            for e in ln.findall(_q(tag)):
                ln.remove(e)
        ln.insert(0, etree.Element(_q('a:noFill')))
    style = sp.find(_q('p:style'))
    if style is not None:
        sp.remove(style)
    tx = sp.find(_q('p:txBody'))
    if tx is not None:
        for p_el in tx.findall(_q('a:p')):
            _make_transparent(p_el)


def _hide_pic(pic):
    blipFill = pic.find(_q('p:blipFill'))
    if blipFill is not None:
        blip = blipFill.find(_q('a:blip'))
        if blip is not None:
            for e in blip.findall(_q('a:alphaModFix')):
                blip.remove(e)
            am = etree.SubElement(blip, _q('a:alphaModFix'))
            am.set('amt', '0')
    spPr = pic.find(_q('p:spPr'))
    if spPr is not None:
        for e in spPr.findall(_q('a:effectLst')):
            spPr.remove(e)
        ln = spPr.find(_q('a:ln'))
        if ln is not None:
            for tag in ('a:solidFill', 'a:gradFill', 'a:pattFill'):
                for e in ln.findall(_q(tag)):
                    ln.remove(e)
            ln.insert(0, etree.Element(_q('a:noFill')))


def _hide_elem(el):
    tag = etree.QName(el).localname
    if tag == 'pic':
        _hide_pic(el)
    elif tag == 'grpSp':
        for ch in el:
            ct = etree.QName(ch).localname
            if ct in ('sp', 'pic', 'grpSp', 'cxnSp'):
                _hide_elem(ch)
    else:  # sp / cxnSp
        _hide_sp_visual(el)


def _has_visible_fill(el):
    """True if a shape draws a filled area: explicit spPr fill, a style
    fill reference (theme fill), or a picture fill."""
    spPr = el.find(_q('p:spPr'))
    if spPr is not None:
        if spPr.find(_q('a:noFill')) is not None:
            pass  # explicit no-fill wins over style refs
        elif any(spPr.find(_q(t)) is not None for t in
                 ('a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill')):
            return True
        else:
            style = el.find(_q('p:style'))
            if style is not None and style.find(_q('a:fillRef')) is not None:
                ref = style.find(_q('a:fillRef'))
                if (ref.get('idx') or '0') not in ('0',):
                    return True
    return el.find(_q('p:blipFill')) is not None


def _has_visible_line(el):
    """True if a shape draws an outline (explicit or via style lnRef)."""
    spPr = el.find(_q('p:spPr'))
    if spPr is not None:
        ln = spPr.find(_q('a:ln'))
        if ln is not None:
            if ln.find(_q('a:noFill')) is not None:
                return False
            if any(ln.find(_q(t)) is not None for t in
                   ('a:solidFill', 'a:gradFill', 'a:pattFill')):
                return True
    style = el.find(_q('p:style'))
    if style is not None:
        ref = style.find(_q('a:lnRef'))
        if ref is not None and (ref.get('idx') or '0') not in ('0',):
            return True
    return False


def _content_elems(root):
    """Direct children of spTree that are drawable elements, in z-order."""
    spTree = root.find(_q('p:cSld') + '/' + _q('p:spTree'))
    out = []
    for ch in spTree:
        t = etree.QName(ch).localname
        if t in ('sp', 'pic', 'grpSp', 'cxnSp'):
            out.append(ch)
    return out


def _elem_bbox(el, slide_w, slide_h):
    tag = etree.QName(el).localname
    pr = el.find(_q('p:grpSpPr')) if tag == 'grpSp' else el.find(_q('p:spPr'))
    if pr is None:
        return None
    xfrm = pr.find(_q('a:xfrm'))
    if xfrm is None:
        return None
    off, ext = xfrm.find(_q('a:off')), xfrm.find(_q('a:ext'))
    if off is None or ext is None:
        return None
    try:
        x, y = int(off.get('x')), int(off.get('y'))
        cx, cy = int(ext.get('cx')), int(ext.get('cy'))
        return (x / slide_w, y / slide_h, (x + cx) / slide_w, (y + cy) / slide_h)
    except (TypeError, ValueError):
        return None


def _shape_pos(sp):
    """(x, y) EMU of a shape, or None if inherited."""
    xfrm = sp.find(_q('p:spPr') + '/' + _q('a:xfrm'))
    if xfrm is None:
        return None
    off = xfrm.find(_q('a:off'))
    if off is None:
        return None
    try:
        return int(off.get('x')), int(off.get('y'))
    except (TypeError, ValueError):
        return None


def _shape_bbox(sp, slide_w, slide_h):
    """Normalized (x0,y0,x1,y1) of a shape, or None."""
    xfrm = sp.find(_q('p:spPr') + '/' + _q('a:xfrm'))
    if xfrm is None:
        return None
    off, ext = xfrm.find(_q('a:off')), xfrm.find(_q('a:ext'))
    if off is None or ext is None:
        return None
    try:
        x, y = int(off.get('x')), int(off.get('y'))
        cx, cy = int(ext.get('cx')), int(ext.get('cy'))
        return (x / slide_w, y / slide_h, (x + cx) / slide_w, (y + cy) / slide_h)
    except (TypeError, ValueError):
        return None


def _is_kpi_text(t):
    """Short numeric statements ('32%', '4.8x', '$1.2M', '250+') are KPIs."""
    t = (t or '').strip()
    if not t or len(t) > 14:
        return False
    digits = sum(c.isdigit() for c in t)
    return digits >= 1 and digits / len(t) >= 0.4


def plan_steps(slide_xml_bytes, slide_w=12192000, slide_h=6858000, script=None):
    """AI Motion Director layout & unit analyzer.
    Detects lines, containers, dots/icons, images, cards, headings, KPIs and
    bullet points, clusters them into semantic narrative units, classifies the
    slide's layout archetype, and returns (groups, bboxes, texts, roles,
    layout) that reveal in narration sequence.
    """
    root = etree.fromstring(slide_xml_bytes)
    elems = _content_elems(root)

    raw_items = []
    for ei, el in enumerate(elems):
        tag = etree.QName(el).localname
        ph = el.findall('.//' + _q('p:ph'))
        if any(p.get('type') in ('title', 'ctrTitle') for p in ph):
            continue
        bb = _elem_bbox(el, slide_w, slide_h)
        if bb is None:
            continue
        x0, y0, x1, y1 = bb
        w, h, area = x1 - x0, y1 - y0, (x1 - x0) * (y1 - y0)

        # Header / branding band: wide title strips AND any element living
        # entirely in the top strip (logos, brand marks) stay static — they
        # must never join a content animation unit
        if (y0 < 0.16 and w > 0.4) or y1 < 0.145:
            continue

        tx = el.find(_q('p:txBody'))
        texts = [_para_text(p) for p in tx.findall(_q('a:p'))] if tx is not None else []
        live = [t for t in texts if t]

        is_image = tag == 'pic'
        is_line = not is_image and ((w > 0.35 and h < 0.05) or tag == 'cxnSp')
        is_dot = not is_image and (w < 0.16 and h < 0.16 and not live)
        has_fill = _has_visible_fill(el)
        is_container = not is_image and ((area > 0.12 and not live and has_fill)
                                         or (tag == 'grpSp' and area > 0.15))

        if len(live) >= 3 and not is_container and not is_image:
            # Split a multi-paragraph body into per-bullet items. Heights are
            # weighted by estimated wrapped-line counts so highlight boxes
            # land on the actual bullet, not an even N-way split.
            cpl = max(14, int(w * 85))            # est. characters per line
            weights = [(max(1.0, math.ceil(len(t) / cpl)) if t else 0.4)
                       for t in texts]
            total = sum(weights) or 1.0
            yc = y0
            for pi, t in enumerate(texts):
                hh = h * weights[pi] / total
                if t:
                    raw_items.append({
                        'kind': 'para', 'ei': ei, 'pi': pi, 'text': t,
                        'bbox': (x0, yc, x1, yc + hh),
                        'is_line': False, 'is_dot': False, 'is_container': False,
                        'is_image': False, 'is_kpi': _is_kpi_text(t),
                    })
                yc += hh
        else:
            raw_items.append({
                'kind': 'elem', 'ei': ei, 'pi': None, 'text': ' '.join(live),
                'bbox': bb, 'is_line': is_line, 'is_dot': is_dot,
                'is_container': is_container, 'is_image': is_image,
                'is_kpi': _is_kpi_text(' '.join(live)) if live else False,
            })

    if not raw_items:
        return [], [], [], [], 'content'

    containers = [it for it in raw_items if it.get('is_container')]
    lines = [it for it in raw_items if it.get('is_line')]
    dots = [it for it in raw_items if it.get('is_dot')]
    content_items = [it for it in raw_items if not it.get('is_container') and not it.get('is_line') and not it.get('is_dot')]
    
    units = []
    used = set()
    
    for ln in lines:
        units.append({'items': [ln], 'type': 'line', 'text': ln['text'] or 'timeline_line'})
        used.add(id(ln))
        
    for c in containers:
        units.append({'items': [c], 'type': 'container', 'text': c['text'] or 'section_container'})
        used.add(id(c))
        
    for d in dots:
        dx0, dy0, dx1, dy1 = d['bbox']
        dcx, dcy = (dx0 + dx1) / 2, (dy0 + dy1) / 2
        matches = []
        for ci in content_items:
            if id(ci) in used: continue
            cx0, cy0, cx1, cy1 = ci['bbox']
            ccx, ccy = (cx0 + cx1) / 2, (cy0 + cy1) / 2
            dist = ((dcx - ccx)**2 + (dcy - ccy)**2)**0.5
            # a dot anchor only absorbs its OWN logo/label: an overlapping
            # image (logo inside the circle) or a small caption hugging the
            # dot — body copy below stays a separate narrative unit so each
            # element animates on its own narration beat
            overlaps = not (cx1 < dx0 - 0.01 or cx0 > dx1 + 0.01 or
                            cy1 < dy0 - 0.01 or cy0 > dy1 + 0.01)
            hugging = (abs(ccx - dcx) < 0.08 and
                       -0.02 < cy0 - dy1 < 0.07)
            small = (cx1 - cx0) * (cy1 - cy0) < 0.03 or ci.get('is_image')
            if small and dist < 0.2 and (overlaps or hugging):
                matches.append((dist, ci))
        matches.sort(key=lambda x: x[0])
        unit_items = [d]
        used.add(id(d))
        for _, ci in matches[:2]:
            unit_items.append(ci)
            used.add(id(ci))
        txt = ' '.join(i['text'] for i in unit_items if i.get('text'))
        units.append({'items': unit_items, 'type': 'marker_or_tool', 'text': txt})
        
    remaining = [it for it in content_items if id(it) not in used]
    headings = [it for it in remaining
                if it['text'] and len(it['text']) < 45
                and not it['text'].endswith('.') and not it.get('is_image')]

    def _heading_like(it):
        t = (it.get('text') or '').strip()
        return bool(t) and len(t) < 45 and not t.endswith('.')

    for h in headings:
        if id(h) in used: continue
        hx0, hy0, hx1, hy1 = h['bbox']
        hcx, hcy = (hx0 + hx1) / 2, (hy0 + hy1) / 2
        best = None
        best_dist = 1e9
        for c in remaining:
            if id(c) in used or c == h: continue
            cx0, cy0, cx1, cy1 = c['bbox']
            ccx, ccy = (cx0 + cx1) / 2, (cy0 + cy1) / 2
            dy_horiz = abs(hcy - ccy)
            dx_horiz = cx0 - hx1
            # Only tight label+bubble pairs merge (a chip touching its text,
            # or a caption on an image). Everything else stays its own unit:
            # the heading pops first, the body pops on ITS narration beat —
            # and two side-by-side headings never fuse into one box.
            is_horiz = (dy_horiz < 0.06 and -0.02 < dx_horiz < 0.05 and
                        (c.get('is_image') or not _heading_like(c)))
            if is_horiz:
                dist = dy_horiz + max(0.0, dx_horiz)
                if dist < best_dist:
                    best_dist = dist
                    best = c
        unit_items = [h]
        used.add(id(h))
        if best:
            unit_items.append(best)
            used.add(id(best))
        txt = ' '.join(i['text'] for i in unit_items if i.get('text'))
        units.append({'items': unit_items,
                      'type': 'card' if best else 'item', 'text': txt})

    for it in remaining:
        if id(it) not in used:
            units.append({'items': [it], 'type': 'item', 'text': it['text']})
            used.add(id(it))
            
    final_units = []
    for u in units:
        txt = u['text'].strip()
        # dot anchors (tool circles, pins) are narrative units even when they
        # carry no text — they pop as the visual anchor of their tool/marker
        if u['type'] in ('line', 'container', 'marker_or_tool') or txt or \
           any(it.get('is_image') for it in u['items']):
            final_units.append(u)

    layout = _classify_layout(final_units)
    for u in final_units:
        u['role'] = _unit_role(u, layout)

    groups = [u['items'] for u in final_units]
    roles = [u['role'] for u in final_units]

    def _boxes(gs):
        out = []
        for g in gs:
            bs = [it['bbox'] for it in g]
            out.append([min(b[0] for b in bs), min(b[1] for b in bs),
                        max(b[2] for b in bs), max(b[3] for b in bs)])
        return out

    gboxes = _boxes(groups)
    gtexts = [' '.join(it.get('text') or '' for it in g).strip() for g in groups]

    # reveal sequence = the slide's natural READING order (never rearranged
    # by the narration): top-to-bottom bands, left column fully before the
    # right column, tool/marker columns cluster circle -> name -> content
    order = _reading_order(gboxes, roles)
    groups = [groups[i] for i in order]
    gboxes = [gboxes[i] for i in order]
    gtexts = [gtexts[i] for i in order]
    roles = [roles[i] for i in order]

    if len(groups) > MAX_STEPS_PER_SLIDE:
        size = -(-len(groups) // MAX_STEPS_PER_SLIDE)

        def _merged_role(rs):
            # a chunk mixing a line/container with real content is content:
            # structural roles would suppress its emphasis & entrance fx
            return next((r for r in rs if r not in ('line', 'container')),
                        rs[0])

        roles = [_merged_role(roles[i:i + size])
                 for i in range(0, len(groups), size)]
        groups = [sum(groups[i:i + size], []) for i in range(0, len(groups), size)]
        gboxes = _boxes(groups)
        gtexts = [' '.join(it.get('text') or '' for it in g).strip() for g in groups]

    return groups, gboxes, gtexts, roles, layout


_TOOL_WORDS = {'chatgpt', 'copilot', 'gemini', 'claude', 'midjourney',
               'perplexity', 'deepseek', 'llama', 'mistral', 'grok',
               'dalle', 'dall-e', 'sora', 'openai', 'anthropic'}


def _unit_role(u, layout):
    """Semantic role of a narrative unit — drives entrance fx & emphasis."""
    t = u['type']
    items = u['items']
    if t == 'line':
        return 'line'
    if t == 'container':
        return 'container'
    if t == 'marker_or_tool':
        if layout == 'timeline':
            return 'marker'
        if layout == 'toolshow' and (
                any(it.get('is_image') for it in items) or len(items) > 1):
            return 'tool'
        return 'card' if len(items) > 1 else 'marker'
    if all(it.get('is_image') for it in items):
        return 'image'
    if all(it.get('is_kpi') for it in items if it.get('text')):
        if any(it.get('is_kpi') for it in items):
            return 'kpi'
    if t == 'card':
        txt = u['text'].strip()
        return 'heading' if len(items) == 1 and len(txt) < 45 else 'card'
    if t == 'item' and items and items[0].get('kind') == 'para':
        return 'bullet'
    return 'item'


def _classify_layout(units):
    """Slide layout archetype from unit geometry — the animation strategy key.
    One of: timeline, toolshow, process, grid, comparison, bullets, visual,
    content."""
    if not units:
        return 'content'

    def ubox(u):
        bs = [it['bbox'] for it in u['items']]
        return (min(b[0] for b in bs), min(b[1] for b in bs),
                max(b[2] for b in bs), max(b[3] for b in bs))

    lines = [u for u in units if u['type'] == 'line']
    anchors = [u for u in units if u['type'] == 'marker_or_tool']
    cards = [u for u in units if u['type'] == 'card']
    paras = [u for u in units
             if u['items'] and u['items'][0].get('kind') == 'para'
             and u['type'] in ('item', 'card')]
    images = [u for u in units if all(it.get('is_image') for it in u['items'])]

    # timeline: a long backbone line + >=2 anchors/cards sitting near it
    for ln in lines:
        lx0, ly0, lx1, ly1 = ubox(ln)
        if lx1 - lx0 < 0.45:
            continue
        lcy = (ly0 + ly1) / 2
        near = [u for u in anchors + cards
                if abs((ubox(u)[1] + ubox(u)[3]) / 2 - lcy) < 0.25]
        if len(near) >= 2:
            return 'timeline'

    # tool showcase: >=2 dot-anchored units with logos/content attached, or
    # a battery of anchors alongside named AI tools
    rich_anchors = [u for u in anchors
                    if len(u['items']) > 1 or
                    any(it.get('is_image') for it in u['items'])]
    if len(rich_anchors) >= 2:
        return 'toolshow'
    tool_names = sum(1 for u in units
                     if any(w in _TOOL_WORDS
                            for w in re.findall(r'[a-z-]+', u['text'].lower())))
    if len(anchors) >= 2 and tool_names >= 2:
        return 'toolshow'

    # dominant image
    if any((ubox(u)[2] - ubox(u)[0]) * (ubox(u)[3] - ubox(u)[1]) > 0.30
           for u in images):
        return 'visual'

    boxes = [ubox(u) for u in units if u['type'] not in ('line', 'container')]
    if len(boxes) == 2:
        (a, b) = boxes
        aw, ah = a[2] - a[0], a[3] - a[1]
        bw, bh = b[2] - b[0], b[3] - b[1]
        if (aw * ah > 0.05 and bw * bh > 0.05 and
                abs((a[1] + a[3]) / 2 - (b[1] + b[3]) / 2) < 0.16 and
                abs((a[0] + a[2]) / 2 - (b[0] + b[2]) / 2) > 0.28):
            return 'comparison'

    if len(boxes) >= 3:
        cys = [(b[1] + b[3]) / 2 for b in boxes]
        cxs = [(b[0] + b[2]) / 2 for b in boxes]
        areas = [max(1e-4, (b[2] - b[0]) * (b[3] - b[1])) for b in boxes]
        one_band = max(cys) - min(cys) < 0.15
        wide_span = max(cxs) - min(cxs) > 0.45
        similar = max(areas) / min(areas) < 4.5
        if one_band and wide_span and similar:
            return 'process'
        # grid: >=4 similar units in >=2 rows and >=2 columns
        if len(boxes) >= 4 and similar:
            def n_clusters(vals, tol):
                cs = []
                for v in sorted(vals):
                    if not cs or v - cs[-1][-1] > tol:
                        cs.append([v])
                    else:
                        cs[-1].append(v)
                return len(cs)
            if n_clusters(cxs, 0.14) >= 2 and n_clusters(cys, 0.12) >= 2:
                return 'grid'

    if len(paras) >= 3:
        pboxes = [ubox(u) for u in paras]
        x0s = [b[0] for b in pboxes]
        if max(x0s) - min(x0s) < 0.10:
            return 'bullets'

    return 'content'


_STOP = set('''a an the and or but if then than so of to in on for with at by from
as is are was were be been being it its this that these those we you they he she
i our your their my me us him her will would can could should may might must do
does did done have has had having not no nor also very just about into over
under out up down again more most other some such only own same too today well
now here there when what which who how why all each both few many'''.split())


def _reading_order(gboxes, roles):
    """The slide's natural reading order, structure-aware.

    Model: content units chain into vertical STACKS (x-overlapping,
    vertically adjacent, inside the same container) — a tool column
    (circle -> name -> description), a bullet list, a heading with its
    text. Stacks group into ROWS by y-overlap. Reading order is rows
    top-to-bottom, blocks left-to-right within a row, and each block
    completes fully before the next — so a left column is read to the
    end before the right column starts.

    Containers reveal first, then backbone lines, then the content.
    Returns the permutation of group indices."""
    n = len(gboxes)
    if n < 2:
        return list(range(n))
    containers = [i for i in range(n) if roles[i] == 'container']
    lines = [i for i in range(n) if roles[i] == 'line']
    rest = [i for i in range(n) if roles[i] not in ('container', 'line')]

    cboxes = [gboxes[i] for i in containers]

    def membership(i):
        b = gboxes[i]
        cx, cy = (b[0] + b[2]) / 2, (b[1] + b[3]) / 2
        return tuple(ci for ci, cb in enumerate(cboxes)
                     if cb[0] - 0.02 <= cx <= cb[2] + 0.02 and
                        cb[1] - 0.02 <= cy <= cb[3] + 0.02)

    def xov(a, b):
        ov = min(a[2], b[2]) - max(a[0], b[0])
        return ov / max(1e-6, min(a[2] - a[0], b[2] - b[0]))

    def stacks_order(idx):
        """Leaf ordering: chain x-aligned adjacent units into stacks
        (a heading + its text, a bullet list), group stacks into rows,
        rows top-to-bottom, stacks left-to-right, units top-to-bottom."""
        stacks = []
        for i in sorted(idx, key=lambda i: (gboxes[i][1], gboxes[i][0])):
            b = gboxes[i]
            wide = (b[2] - b[0]) > 0.5
            best = None
            if not wide:
                # join the NEAREST eligible stack above (greatest bottom) —
                # joining the first-created one can leapfrog an overlapping
                # middle unit and scramble the column's internal order
                for st in stacks:
                    if st['wide'] or st['mem'] != membership(i):
                        continue
                    if xov(st['bbox'], b) >= 0.5 and \
                            -0.06 < b[1] - st['bbox'][3] < 0.16:
                        if best is None or st['bbox'][3] > best['bbox'][3]:
                            best = st
            if best is not None:
                best['idx'].append(i)
                best['bbox'] = [min(best['bbox'][0], b[0]),
                                min(best['bbox'][1], b[1]),
                                max(best['bbox'][2], b[2]),
                                max(best['bbox'][3], b[3])]
            else:
                stacks.append({'idx': [i], 'bbox': list(b),
                               'mem': membership(i), 'wide': wide})
        rows = []
        for st in sorted(stacks, key=lambda s: (s['bbox'][1], s['bbox'][0])):
            b = st['bbox']
            placed = False
            for row in rows:
                ov = min(row['y1'], b[3]) - max(row['y0'], b[1])
                if ov / max(1e-6, min(row['y1'] - row['y0'],
                                      b[3] - b[1])) >= 0.5:
                    row['blocks'].append(st)
                    row['y0'] = min(row['y0'], b[1])
                    row['y1'] = max(row['y1'], b[3])
                    placed = True
                    break
            if not placed:
                rows.append({'blocks': [st], 'y0': b[1], 'y1': b[3]})
        res = []
        for row in rows:
            for st in sorted(row['blocks'], key=lambda s: s['bbox'][0]):
                res.extend(sorted(st['idx'],
                                  key=lambda i: (gboxes[i][1], gboxes[i][0])))
        return res

    # recursive XY-cut: split at clear x-gaps (panels, left before right)
    # then y-gaps (bands, top before bottom); leaves use stack ordering.
    GAP = 0.035

    def cut(idx, axis, depth):
        if len(idx) <= 2 or depth > 8:
            return stacks_order(idx)
        lo, hi = (0, 2) if axis == 'x' else (1, 3)
        ivs = sorted((gboxes[i][lo], gboxes[i][hi], i) for i in idx)
        bands = [[ivs[0][2]]]
        cur_end = ivs[0][1]
        for a, b_, i in ivs[1:]:
            if a - cur_end >= GAP:
                bands.append([i])
            else:
                bands[-1].append(i)
            cur_end = max(cur_end, b_)
        nxt = 'y' if axis == 'x' else 'x'
        if len(bands) == 1:
            if depth == 0 or axis == 'x':
                return cut(idx, nxt, depth + 1)
            return stacks_order(idx)
        out = []
        for band in bands:
            out.extend(cut(band, nxt, depth + 1))
        return out

    out = sorted(containers, key=lambda i: (gboxes[i][1], gboxes[i][0]))
    out += sorted(lines, key=lambda i: (gboxes[i][1], gboxes[i][0]))
    out += cut(rest, 'x', 0)

    # safety: anything containing another unit still precedes it
    def contains(a, b):
        return (a[0] <= b[0] + 0.02 and a[1] <= b[1] + 0.02 and
                a[2] >= b[2] - 0.02 and a[3] >= b[3] - 0.02 and
                (a[2] - a[0]) * (a[3] - a[1]) >
                (b[2] - b[0]) * (b[3] - b[1]) * 1.3)

    guard = 0
    changed = True
    while changed and guard < 60:
        changed = False
        guard += 1
        for j in range(len(out)):
            for i in range(j):
                if contains(gboxes[out[j]], gboxes[out[i]]):
                    out.insert(i, out.pop(j))
                    changed = True
                    break
            if changed:
                break
    return out


def build_states_pptx(src_pptx, out_pptx, scripts=None):
    """Create a pptx whose slide list is every build state of every slide,
    in order. Returns manifest: list per original slide of state counts &
    revealed-paragraph plan."""
    zin = zipfile.ZipFile(src_pptx)
    names = zin.namelist()
    pres_xml = zin.read('ppt/presentation.xml')
    pres_root = etree.fromstring(pres_xml)
    rels_xml = zin.read('ppt/_rels/presentation.xml.rels')
    rels_root = etree.fromstring(rels_xml)
    ct_xml = zin.read('[Content_Types].xml')
    ct_root = etree.fromstring(ct_xml)

    # slide dimensions (for header-band detection)
    sldSz = pres_root.find(_q('p:sldSz'))
    slide_w = int(sldSz.get('cx')) if sldSz is not None else 12192000
    slide_h = int(sldSz.get('cy')) if sldSz is not None else 6858000

    # ordered original slide part names
    sldIdLst = pres_root.find(_q('p:sldIdLst'))
    rid_to_target = {rel.get('Id'): rel.get('Target')
                     for rel in rels_root}
    slide_parts = []
    for sldId in sldIdLst:
        rid = sldId.get(_q('r:id'))
        tgt = rid_to_target[rid]
        slide_parts.append('ppt/' + tgt.lstrip('/').replace('../', ''))

    manifest = []
    new_parts = {}          # part name -> bytes
    new_rels_entries = []   # (rid, target)
    new_sldids = []
    rid_counter = 9000
    sld_counter = 90000
    part_counter = 1

    for orig_idx, part in enumerate(slide_parts):
        slide_bytes = zin.read(part)
        script = (scripts[orig_idx] if scripts and orig_idx < len(scripts)
                  else None)
        groups, gboxes, gtexts, roles, layout = plan_steps(
            slide_bytes, slide_w, slide_h, script=script)
        n_states = len(groups) + 1  # state 0 .. len(groups)
        state_plan = []
        for state in range(n_states):
            root = etree.fromstring(slide_bytes)
            elems = _content_elems(root)
            # Hide only truly-empty PLACEHOLDER shapes (layout prompt boxes).
            # Never touch regular shapes: an outline-only arrow, a divider
            # line, or a theme-filled box IS slide content and must render.
            for ei, el in enumerate(elems):
                tag = etree.QName(el).localname
                if tag == 'sp':
                    ph = el.findall('.//' + _q('p:ph'))
                    tx = el.find(_q('p:txBody'))
                    texts = [_para_text(p) for p in tx.findall(_q('a:p'))] if tx is not None else []
                    live = [t for t in texts if t]
                    if (not live and ph and not _has_visible_fill(el)
                            and not _has_visible_line(el)):
                        _hide_sp_visual(el)
            hidden = []
            for gi, group in enumerate(groups):
                if gi >= state:           # not yet revealed
                    for it in group:
                        try:
                            el = elems[it['ei']]
                            if it['kind'] == 'para':
                                tx = el.find(_q('p:txBody'))
                                _make_transparent(tx.findall(_q('a:p'))[it['pi']])
                            else:
                                _hide_elem(el)
                            hidden.append([it['ei'], it.get('pi')])
                        except Exception:
                            pass
            new_name = f'ppt/slides/slideBS{part_counter}.xml'
            new_parts[new_name] = etree.tostring(root, xml_declaration=True,
                                                 encoding='UTF-8', standalone=True)
            # rels for the new part = copy of original slide rels
            orig_rels = part.replace('ppt/slides/', 'ppt/slides/_rels/') + '.rels'
            if orig_rels in names:
                new_parts[f'ppt/slides/_rels/slideBS{part_counter}.xml.rels'] = \
                    zin.read(orig_rels)
            rid = f'rIdBS{rid_counter}'; rid_counter += 1
            new_rels_entries.append((rid, f'slides/slideBS{part_counter}.xml'))
            new_sldids.append((sld_counter, rid)); sld_counter += 1
            part_counter += 1
            state_plan.append(hidden)
        manifest.append({'orig_index': orig_idx, 'n_states': n_states,
                         'regions': gboxes, 'texts': gtexts,
                         'gsizes': [len(g) for g in groups],
                         'roles': roles, 'layout': layout,
                         'bver': BUILD_VERSION,
                         'aspect': round(slide_w / max(1, slide_h), 4)})

    # rewrite sldIdLst
    for child in list(sldIdLst):
        sldIdLst.remove(child)
    for (sid, rid) in new_sldids:
        el = etree.SubElement(sldIdLst, _q('p:sldId'))
        el.set('id', str(sid)); el.set(_q('r:id'), rid)
    # add relationships
    for (rid, tgt) in new_rels_entries:
        rel = etree.SubElement(rels_root, _q('rel:Relationship'))
        rel.set('Id', rid)
        rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
        rel.set('Target', tgt)
    # content types: ensure slide override for each new part
    for name in new_parts:
        if name.endswith('.rels'):
            continue
        ov = etree.SubElement(ct_root, _q('ct:Override'))
        ov.set('PartName', '/' + name)
        ov.set('ContentType',
               'application/vnd.openxmlformats-officedocument.presentationml.slide+xml')

    with zipfile.ZipFile(out_pptx, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name in names:
            if name == 'ppt/presentation.xml':
                zout.writestr(name, etree.tostring(pres_root, xml_declaration=True,
                                                   encoding='UTF-8', standalone=True))
            elif name == 'ppt/_rels/presentation.xml.rels':
                zout.writestr(name, etree.tostring(rels_root, xml_declaration=True,
                                                   encoding='UTF-8', standalone=True))
            elif name == '[Content_Types].xml':
                zout.writestr(name, etree.tostring(ct_root, xml_declaration=True,
                                                   encoding='UTF-8', standalone=True))
            else:
                zout.writestr(name, zin.read(name))
        for name, data in new_parts.items():
            zout.writestr(name, data)
    zin.close()
    return manifest


def render_build_states(src_pptx, workdir, dpi=110, scripts=None):
    """Full pipeline: returns per-slide list of state image paths + manifest.
    scripts: optional per-slide narration used to order reveals by story."""
    workdir = os.path.abspath(workdir)
    os.makedirs(workdir, exist_ok=True)
    bs_pptx = os.path.join(workdir, 'build_states.pptx')
    manifest = build_states_pptx(src_pptx, bs_pptx, scripts=scripts)
    expected = sum(m['n_states'] for m in manifest)
    # export into a fresh dir, swap in only on success: a partial export must
    # never destroy the previous good states or misalign the state mapping
    tmpdir = os.path.join(workdir, 'states_new')
    if os.path.isdir(tmpdir):
        shutil.rmtree(tmpdir)
    pngs = pptx_to_pngs(bs_pptx, tmpdir, dpi=dpi, prefix='st', expected=expected)
    if len(pngs) != expected:
        raise RuntimeError(f'Build-state export mismatch: expected {expected} '
                           f'state images, got {len(pngs)}')
    # Swap via rename-aside: the live 'states' dir (served to the editor,
    # possibly held open by OneDrive/antivirus) is never deleted before its
    # replacement is in place — a locked file can abort the swap but can
    # never destroy the previous good states.
    imgdir = os.path.join(workdir, 'states')
    olddir = os.path.join(workdir, 'states_old')

    def _retry(fn, *a):
        for k in range(5):
            try:
                return fn(*a)
            except OSError:
                if k == 4:
                    raise
                import time; time.sleep(0.5)

    shutil.rmtree(olddir, ignore_errors=True)
    if os.path.isdir(imgdir):
        _retry(os.replace, imgdir, olddir)
    try:
        _retry(os.replace, tmpdir, imgdir)
    except OSError:
        if os.path.isdir(olddir):     # restore the previous good states
            _retry(os.replace, olddir, imgdir)
        raise
    shutil.rmtree(olddir, ignore_errors=True)
    pngs = [p.replace(tmpdir, imgdir) for p in pngs]
    out, cursor = [], 0
    for m in manifest:
        out.append(pngs[cursor:cursor + m['n_states']])
        cursor += m['n_states']
    return out, manifest
